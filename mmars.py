import streamlit as st
import os
import io
from pandas import pandas as pd
from PyPDF2 import PdfReader
from pptx import Presentation
import docx
import xlrd
from bs4 import BeautifulSoup
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.embeddings import GooglePalmEmbeddings
from langchain_community.llms import GooglePalm
from langchain_community.vectorstores import FAISS
from langchain.chains import ConversationalRetrievalChain
from langchain.memory import ConversationBufferMemory
import cv2
import pytesseract
import numpy as np
from langchain_core.messages import HumanMessage, AIMessage
import logging
from typing import Union
import speech_recognition as sr
from datetime import datetime
import spacy
from config import GOOGLE_API_KEY
import zipfile
import rarfile
import markdown
from textract import process
from transformers import pipeline
from langchain_community.chat_models import ChatOllama

# Initialize NLP model
nlp = spacy.load("en_core_web_sm")

# Set up environment variables for Google API key
os.environ['GOOGLE_API_KEY'] = GOOGLE_API_KEY

# Create a logger
logger = logging.getLogger(__name__)

# Mapping file extensions to their respective text extraction functions
EXTRACTION_FUNCTIONS = {
    "pdf": "extract_text_from_pdf",
    "pptx": "extract_text_from_ppt",
    "docx": "extract_text_from_docx",
    "py": "extract_text_from_py",
    "xls": "extract_text_from_excel",
    "xlsx": "extract_text_from_excel",
    "csv": "extract_text_from_csv",
    "html": "extract_text_from_html",
    "css": "extract_text_from_css",
    "json": "extract_text_from_json",
    "sql": "extract_text_from_sql",
    "txt": "extract_text_from_txt",
    "java": "extract_text_from_java",
    "c": "extract_text_from_c",
    "cpp": "extract_text_from_cpp",
    "js": "extract_text_from_javascript",
    "swift": "extract_text_from_swift",
    "r": "extract_text_from_r",
    "rs": "extract_text_from_rust",
    "jpg": "extract_text_from_image",
    "jpeg": "extract_text_from_image",
    "png": "extract_text_from_image",
    "bmp": "extract_text_from_image",
    "xml": "extract_text_from_xml",
    "md": "extract_text_from_md",
    "tex": "extract_text_from_tex",
    "zip": "extract_text_from_zip",
    "rar": "extract_text_from_rar",
}

def extract_text(file):
    file_extension = file.name.split(".")[-1].lower()
    extraction_function = EXTRACTION_FUNCTIONS.get(file_extension)

    if extraction_function:
        return globals()[extraction_function](file)
    else:
        handle_file_processing_error(file_extension, "Unsupported file type")
        return ""

def extract_text_from_pdf(pdf_file):
    return extract_from_pdf_or_doc(pdf_file, "pdf")

def extract_text_from_ppt(ppt_file):
    return extract_from_ppt_or_docx(ppt_file, "pptx")

def extract_text_from_docx(docx_file):
    return extract_from_ppt_or_docx(docx_file, "docx")

def extract_from_pdf_or_doc(file, file_type):
    text = ""
    try:
        if file_type == "pdf":
            reader = PdfReader(file)
            for page in reader.pages:
                text += page.extract_text()
        else:
            doc = docx.Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
    except Exception as e:
        handle_file_processing_error(file_type, e)
    return text

def extract_from_ppt_or_docx(file, file_type):
    text = ""
    try:
        if file_type == "pptx":
            presentation = Presentation(file)
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text
        else:
            doc = docx.Document(file)
            for para in doc.paragraphs:
                text += para.text + "\n"
    except Exception as e:
        handle_file_processing_error(file_type, e)
    return text

def handle_file_processing_error(file_type: str, error: Exception):
    st.error(f"Error processing {file_type} file: {error}")
    logger.exception(f"Error processing {file_type} file", exc_info=True)

def main():
    st.set_page_config(page_title="M.A.R.S 🚀", layout="wide")
    st.header("Multi-modal AI Research System")
    user_question = st.chat_input("Ask Questions about Everything")

    # Handle file upload and text extraction
    files = st.file_uploader("Upload Files", accept_multiple_files=True)
    if files:
        raw_text = ""
        for file in files:
            raw_text += extract_text(file)
        
        # Process text, vectorize, and handle AI model interaction
        if raw_text:
            st.success("Files processed successfully")
            # Further processing and AI interaction logic here
        else:
            st.warning("No text extracted from files.")

if __name__ == "__main__":
    main()
